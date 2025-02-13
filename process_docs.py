
#import os
#import pdfminer.high_level
##from langchain.embeddings import HuggingFaceEmbeddings
#from langchain_ollama import OllamaEmbeddings
#import chromadb
#
## Directory where your documents are stored
#DOCUMENTS_DIR = "/home/ai-bench/Documents"
#
## Initialize ChromaDB
#chroma_client = chromadb.PersistentClient(path="./chroma_db")
#collection = chroma_client.get_or_create_collection(name="device_docs")
#
## Load embedding model
##embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
#embedding_model = OllamaEmbeddings(model="deepseek-r1")
## Function to extract text from PDF
#def extract_text_from_pdf(pdf_path):
#    return pdfminer.high_level.extract_text(pdf_path)
#
## Process all PDFs in the directory
#for filename in os.listdir(DOCUMENTS_DIR):
#    if filename.endswith(".pdf"):
#        file_path = os.path.join(DOCUMENTS_DIR, filename)
#        text = extract_text_from_pdf(file_path)
#        
#        # Split text into smaller chunks (e.g., 500-character chunks)
#        chunks = [text[i:i+500] for i in range(0, len(text), 500)]
#        
#        # Add to ChromaDB
#        collection.add(
#            documents=chunks,
#            metadatas=[{"source": filename}] * len(chunks),
#            ids=[f"{filename}-{i}" for i in range(len(chunks))]
#        )
#
#print("Documents successfully indexed!")
import os
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import pdfminer.high_level
from pptx import Presentation
import docx
import chromadb
from langchain_ollama import OllamaEmbeddings

# Directory where your documents are stored
DOCUMENTS_DIR = "/home/ai-bench/Documents"

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(path="./chroma_db")
collection = chroma_client.get_or_create_collection(name="device_docs")

# Load embedding model
embedding_model = OllamaEmbeddings(model="deepseek-r1")

# Function to extract text from PDFs
def extract_text_from_pdf(pdf_path):
    return pdfminer.high_level.extract_text(pdf_path)

# Function to extract text from PowerPoint files
def extract_text_from_pptx(ppt_path):
    prs = Presentation(ppt_path)
    text = "\n".join([slide.notes_slide.text_frame.text for slide in prs.slides if slide.notes_slide])
    return text or "No text extracted."

# Function to extract text from Word documents
def extract_text_from_docx(doc_path):
    doc = docx.Document(doc_path)
    return "\n".join([para.text for para in doc.paragraphs])

# Function to extract text from plain text files
def extract_text_from_txt(txt_path):
    with open(txt_path, "r", encoding="utf-8") as f:
        return f.read()

# Function to process a new document
def process_document(file_path):
    _, ext = os.path.splitext(file_path)
    text = ""

    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext == ".txt":
        text = extract_text_from_txt(file_path)

    if text:
        chunks = [text[i:i+500] for i in range(0, len(text), 500)]
        collection.add(
            documents=chunks,
            metadatas=[{"source": os.path.basename(file_path)}] * len(chunks),
            ids=[f"{os.path.basename(file_path)}-{i}" for i in range(len(chunks))]
        )
        print(f"Indexed: {file_path}")

# ** Watchdog Handler for Monitoring New Files **
class DocumentHandler(FileSystemEventHandler):
    def on_created(self, event):
        if not event.is_directory:
            process_document(event.src_path)

# Start Watching for New Files
if __name__ == "__main__":
    event_handler = DocumentHandler()
    observer = Observer()
    observer.schedule(event_handler, DOCUMENTS_DIR, recursive=False)
    observer.start()

    print("Watching for new documents in:", DOCUMENTS_DIR)
    try:
        while True:
            time.sleep(5)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()

